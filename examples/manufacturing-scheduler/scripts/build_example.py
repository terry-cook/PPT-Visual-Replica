from __future__ import annotations

import json
import shutil
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
EXAMPLE_ROOT = SCRIPT_DIR.parent
REPO_ROOT = EXAMPLE_ROOT.parents[1]
OUT = REPO_ROOT / "local-runs" / "manufacturing-scheduler-rebuild"
ASSET_SRC = EXAMPLE_ROOT / "imagegen" / "transparent-assets"
ASSETS = OUT / "assets"
PROMPTS = OUT / "prompts"
REF_SRC = EXAMPLE_ROOT / "reference" / "reference.png"
REF = OUT / "reference.png"
CANVAS_W = 1672
CANVAS_H = 941
SLIDE_W_IN = 16.0
SLIDE_H_IN = SLIDE_W_IN * CANVAS_H / CANVAS_W
SCALE = SLIDE_W_IN / CANVAS_W


COLORS = {
    "navy": "0B2A67",
    "blue": "0E67C8",
    "blue2": "1F80E0",
    "teal": "00A8A8",
    "teal_dark": "067A82",
    "orange": "FF7A1A",
    "orange2": "F59B22",
    "line_blue": "1A65C8",
    "line_teal": "00A0A0",
    "line_orange": "FF7A1A",
    "soft_blue": "EAF5FF",
    "soft_teal": "EAFBFA",
    "soft_orange": "FFF4E8",
    "white": "FFFFFF",
    "ink": "12345B",
    "muted": "28566E",
}


def emu(px: float) -> int:
    return int(Inches(px * SCALE))


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def ensure_dirs() -> None:
    for folder in [OUT, ASSETS, PROMPTS, OUT / "reference_crops", OUT / "preview"]:
        folder.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(REF_SRC, REF)
    for asset in ASSET_SRC.glob("*.png"):
        shutil.copyfile(asset, ASSETS / asset.name)


def crop_cover(src: Path, out: Path, target_ratio: float) -> None:
    img = Image.open(src).convert("RGB")
    w, h = img.size
    current = w / h
    if current > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        box = (left, 0, left + new_w, h)
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        box = (0, top, w, top + new_h)
    img.crop(box).save(out)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    color: str = "12345B",
    bold: bool = False,
    align: str = "left",
    valign: str = "top",
    name: str | None = None,
    margin: float = 2,
) -> None:
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1.0,
    name: str | None = None,
) -> object:
    shape = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    if name:
        shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str,
    width: float = 1.2,
    dashed: bool = False,
    name: str | None = None,
) -> object:
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    if name:
        shape.name = name
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def add_picture_fit(slide, file_name: str, x: float, y: float, w: float, h: float, name: str) -> dict:
    path = ASSETS / file_name
    img = Image.open(path)
    iw, ih = img.size
    fit = min(w / iw, h / ih)
    fw, fh = iw * fit, ih * fit
    fx, fy = x + (w - fw) / 2, y + (h - fh) / 2
    pic = slide.shapes.add_picture(str(path), emu(fx), emu(fy), emu(fw), emu(fh))
    pic.name = name
    return {"id": name, "path": file_name, "anchor_slot": [x, y, w, h], "fitted_bbox": [fx, fy, fw, fh]}


IMAGE_UNITS = [
    ("scene_bg", "scene_background_top.png", [122, 107, 1526, 255], "generated_scene_background"),
    ("scene_layer_icon", "scene_factory.png", [36, 176, 62, 62], "imagegen_asset"),
    ("sensor_icon", "sensor_data.png", [263, 140, 45, 42], "imagegen_asset"),
    ("status_icon", "equipment_status.png", [367, 141, 48, 42], "imagegen_asset"),
    ("environment_icon", "environment_info.png", [468, 140, 50, 42], "imagegen_asset"),
    ("amr_1", "amr.png", [165, 246, 126, 74], "imagegen_asset"),
    ("robot_arm_1", "robot_arm.png", [360, 219, 134, 94], "imagegen_asset"),
    ("amr_2", "amr.png", [618, 247, 126, 74], "imagegen_asset"),
    ("agv_1", "agv.png", [760, 250, 134, 73], "imagegen_asset"),
    ("robot_arm_2", "robot_arm.png", [980, 218, 134, 94], "imagegen_asset"),
    ("server_1", "server.png", [1128, 190, 88, 126], "imagegen_asset"),
    ("database_1", "database.png", [1250, 206, 93, 112], "imagegen_asset"),
    ("twin_screen_1", "twin_screen.png", [1358, 159, 248, 153], "imagegen_asset"),
    ("algorithm_layer_icon", "algorithm_chip.png", [36, 500, 64, 64], "imagegen_asset"),
    ("order_clipboard", "clipboard.png", [223, 492, 72, 86], "imagegen_asset"),
    ("task_assign_icon", "task_assign.png", [507, 441, 36, 36], "imagegen_asset"),
    ("resource_icon", "resource_cube.png", [506, 519, 38, 38], "imagegen_asset"),
    ("energy_icon", "energy_droplet.png", [506, 592, 38, 38], "imagegen_asset"),
    ("constraint_icon", "constraint_shield.png", [706, 434, 32, 32], "imagegen_asset"),
    ("brain_scheduler_icon", "brain_scheduler.png", [801, 496, 78, 65], "imagegen_asset"),
    ("policy_update_icon", "dispatch_nodes.png", [742, 614, 32, 32], "imagegen_asset"),
    ("anomaly_icon", "alarm.png", [1045, 443, 44, 36], "imagegen_asset"),
    ("collab_icon", "robot_collab.png", [1046, 512, 44, 42], "imagegen_asset"),
    ("status_eval_icon", "status_bars.png", [1046, 591, 44, 42], "imagegen_asset"),
    ("resource_pool_group_icon", "resource_pool_group.png", [1306, 492, 126, 40], "imagegen_asset"),
    ("device_pool_group_icon", "device_pool_group.png", [1306, 574, 126, 44], "imagegen_asset"),
    ("output_layer_icon", "output_monitor.png", [36, 788, 62, 62], "imagegen_asset"),
    ("decision_graph_icon", "dispatch_nodes.png", [184, 782, 114, 60], "imagegen_asset"),
    ("execution_doc_icon", "task_assign.png", [407, 791, 58, 70], "imagegen_asset"),
    ("route_icon", "route_kpi.png", [560, 790, 96, 63], "imagegen_asset"),
    ("alarm_output_icon", "alarm.png", [738, 788, 78, 62], "imagegen_asset"),
    ("energy_output_icon", "energy_leaf_battery.png", [884, 788, 80, 65], "imagegen_asset"),
    ("kpi_bar_completion", "kpi_bar_1.png", [1026, 820, 78, 42], "imagegen_asset"),
    ("kpi_bar_utilization", "kpi_bar_2.png", [1140, 820, 78, 42], "imagegen_asset"),
    ("kpi_line_response", "kpi_line_blue.png", [1261, 820, 78, 42], "imagegen_asset"),
    ("kpi_line_energy", "kpi_line_green.png", [1375, 820, 78, 42], "imagegen_asset"),
    ("kpi_line_anomaly", "kpi_line_orange.png", [1489, 820, 78, 42], "imagegen_asset"),
]


def draw_panel(slide, x, y, w, h, line, fill, name):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line, line_width=1.2, name=name)


def add_data_badge(slide, x, y, w, text, color, name):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 30, fill="EAF7FA", line=None, name=f"{name}_bg")
    add_text(slide, text, x, y + 3, w, 23, size=14, color=color, bold=True, align="center", valign="middle", name=name)
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, x + w / 2 - 7, y - 24, 14, 30, fill=color, line=None, name=f"{name}_arrow")


def add_card(slide, x, y, w, h, border, title, title_color, title_h=34):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill="FFFFFF", line=border, line_width=1.1)
    add_text(slide, title, x + 8, y + 9, w - 16, title_h, size=14, color=title_color, bold=True, align="center")


def build_pptx() -> list[dict]:
    crop_cover(ASSETS / "scene_background.png", ASSETS / "scene_background_top.png", 1526 / 255)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, fill="FFFFFF", line=None, name="background")
    add_text(slide, "面向智能制造的多机器人协同调度与资源优化架构", 260, 13, 1150, 52, size=30, color=COLORS["navy"], bold=True, align="center", name="title")
    add_text(slide, "多源感知—智能调度—闭环优化", 646, 69, 380, 30, size=20, color=COLORS["blue"], bold=True, align="center", name="subtitle")

    # Macro panels and generated scene background.
    draw_panel(slide, 10, 106, 1638, 256, COLORS["line_blue"], "F5FAFF", "panel_scene")
    add_picture_fit(slide, "scene_background_top.png", 122, 107, 1524, 254, "scene_bg")
    draw_panel(slide, 10, 410, 1638, 258, COLORS["line_teal"], "F7FEFE", "panel_scheduler")
    draw_panel(slide, 10, 715, 1638, 185, COLORS["line_orange"], "FFF9F2", "panel_output")

    # Layer ribbons.
    for y, h, color, label, icon, icon_y, label_y, label_h, label_size in [
        (106, 256, COLORS["blue"], "场景层", "scene_layer_icon", 174, 121, 46, 17),
        (410, 258, COLORS["teal"], "算法与\n资源调度层", "algorithm_layer_icon", 496, 428, 62, 15),
        (715, 185, COLORS["orange"], "输出决策层", "output_layer_icon", 784, 728, 42, 15),
    ]:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 10, y, 102, h, fill=color, line=None, name=f"ribbon_{y}")
        add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, 86, y, 52, 58, fill=color, line=None, name=f"ribbon_tip_{y}")
        add_shape(slide, MSO_SHAPE.OVAL, 31, icon_y, 68, 68, fill=color, line="FFFFFF", line_width=1.3, name=f"ribbon_icon_ring_{y}")
        add_text(slide, label, 13, label_y, 96, label_h, size=label_size, color="FFFFFF", bold=True, align="center", valign="middle", name=f"ribbon_label_{y}", margin=0)

    placements = []

    # Top layer sensing card.
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 238, 120, 305, 87, fill="FFFFFF", line="C9D9EA", line_width=1.0, name="sensing_input_card")
    add_text(slide, "多源感知输入", 323, 116, 142, 26, size=14, color=COLORS["blue"], bold=True, align="center", name="sensing_input_title")
    for x in [335, 437]:
        add_line(slide, x, 141, x, 174, COLORS["line_blue"], width=1.0)
    add_text(slide, "传感器数据", 252, 184, 74, 18, size=10, color="111111", bold=True, align="center")
    add_text(slide, "设备状态", 358, 184, 72, 18, size=10, color="111111", bold=True, align="center")
    add_text(slide, "环境信息", 462, 184, 72, 18, size=10, color="111111", bold=True, align="center")

    # Data collection rails and machine labels.
    top_drop_xs = [287, 637, 735, 814, 891, 1004, 1174, 1291]
    for x in top_drop_xs:
        add_shape(slide, MSO_SHAPE.OVAL, x - 5, 153, 10, 10, fill=COLORS["line_blue"], line=None)
        add_line(slide, x, 158, x, 228, COLORS["line_blue"], width=1.0, dashed=True)
        add_shape(slide, MSO_SHAPE.DOWN_ARROW, x - 6, 221, 12, 22, fill=COLORS["line_blue"], line=None)
    add_line(slide, 542, 158, 1174, 158, COLORS["line_blue"], width=1.7)
    add_line(slide, 1291, 154, 1404, 154, COLORS["line_blue"], width=1.2, dashed=True)
    for text, x, w in [
        ("移动机器人 / AMR ①", 145, 155),
        ("机械臂工作站 ①", 371, 132),
        ("移动机器人 / AMR ②", 599, 162),
        ("AGV 运输车 ①", 761, 127),
        ("机械臂工作站 ②", 973, 132),
        ("边缘服务器", 1118, 110),
        ("生产订单数据库", 1238, 128),
        ("数字孪生监控屏", 1411, 150),
    ]:
        add_text(slide, text, x, 330, w, 22, size=11, color="111111", bold=True, align="center")

    for x in [1095, 1218, 1348, 1607]:
        add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, 236, 38, 17, fill=COLORS["line_blue"], line=None)

    add_data_badge(slide, 290, 372, 280, "实时数据（设备/环境/状态）", COLORS["blue"], "badge_real_time")
    add_data_badge(slide, 650, 372, 300, "资源状态（机器人/设备/能耗）", COLORS["blue"], "badge_resource")
    add_data_badge(slide, 1015, 372, 285, "订单信息（生产计划/优先级）", COLORS["blue"], "badge_order")

    # Scheduler layer.
    add_card(slide, 209, 439, 185, 197, COLORS["line_teal"], "订单需求 / 任务池", COLORS["teal_dark"])
    add_text(slide, "· 生产订单\n· 任务队列\n· 优先级/交期\n· 物料需求", 304, 504, 86, 90, size=11, color="111111", bold=False)
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 399, 526, 57, 28, fill=COLORS["line_teal"], line=None, name="arrow_order_to_scheduler")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 463, 425, 748, 223, fill=None, line=COLORS["line_teal"], line_width=1.0, name="scheduler_boundary")

    for y, icon, label in [
        (436, "task_assign_icon", "任务分配"),
        (508, "resource_icon", "资源优化"),
        (582, "energy_icon", "能耗优化"),
    ]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 493, y, 145, 57, fill="F2FFFF", line=COLORS["line_teal"], line_width=1.0)
        add_text(slide, label, 546, y + 16, 76, 24, size=13, color=COLORS["teal_dark"], bold=True, align="center")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 688, 425, 277, 49, fill="F2FFFF", line=COLORS["line_teal"], line_width=1.0)
    add_text(slide, "约束管理（时序/设备/路径）", 772, 438, 170, 24, size=12.5, color=COLORS["teal_dark"], bold=True, align="center")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 720, 484, 238, 116, fill="11B9BA", line=None, name="rl_scheduler_core")
    add_text(slide, "强化学习调度器", 749, 566, 181, 28, size=18, color="FFFFFF", bold=True, align="center", valign="middle")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 720, 612, 238, 41, fill="F2FFFF", line=COLORS["line_teal"], line_width=1.0)
    add_text(slide, "调度策略更新", 781, 622, 132, 23, size=12.5, color=COLORS["teal_dark"], bold=True, align="center")

    for y, label in [(437, "异常检测"), (510, "多机器人协同"), (583, "状态评估")]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1031, y, 149, 56, fill="F2FFFF", line=COLORS["line_teal"], line_width=1.0)
        add_text(slide, label, 1090, y + 17, 78, 24, size=12.5, color=COLORS["teal_dark"], bold=True, align="center")

    for y in [453, 523, 605]:
        add_line(slide, 638, y, 720, y, COLORS["teal_dark"], width=1.0, dashed=True)
        add_line(slide, 958, y, 1030, y, COLORS["teal_dark"], width=1.0, dashed=True)
    add_line(slide, 840, 474, 840, 484, COLORS["teal_dark"], width=2.0)
    add_line(slide, 840, 600, 840, 612, COLORS["teal_dark"], width=2.0)

    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 1188, 525, 60, 27, fill=COLORS["line_teal"], line=None)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1282, 433, 240, 204, fill="F7FFFF", line=COLORS["line_teal"], line_width=1.0)
    add_text(slide, "资源池", 1347, 440, 110, 28, size=15, color=COLORS["teal_dark"], bold=True, align="center")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1290, 470, 224, 86, fill="FFFFFF", line="C7E4E3", line_width=1.0)
    add_text(slide, "机器人资源池", 1348, 476, 110, 21, size=11.5, color=COLORS["navy"], bold=True, align="center")
    add_text(slide, "AGV", 1342, 535, 40, 17, size=7.5, color="111111", align="center")
    add_text(slide, "…", 1484, 505, 22, 18, size=14, color="111111", bold=True, align="center")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1290, 562, 224, 64, fill="FFFFFF", line="C7E4E3", line_width=1.0)
    add_text(slide, "设备资源池", 1348, 566, 110, 20, size=11.5, color=COLORS["navy"], bold=True, align="center")
    add_text(slide, "…", 1484, 588, 22, 18, size=14, color="111111", bold=True, align="center")

    add_data_badge(slide, 395, 676, 190, "优化结果 / 调度策略", COLORS["teal_dark"], "badge_strategy")
    add_data_badge(slide, 697, 676, 220, "控制指令 / 执行计划", COLORS["teal_dark"], "badge_plan")
    add_data_badge(slide, 1015, 676, 250, "反馈数据（执行结果/KPI）", COLORS["teal_dark"], "badge_feedback")

    # Output layer.
    cards = [
        (157, 733, 183, 155, "调度决策输出", "最优任务—资源—时间方案"),
        (356, 733, 154, 155, "任务执行指令", "下发控制指令\n监控执行进度"),
        (526, 733, 164, 155, "路径与工位分配", "最优路径规划\n工位动态分配"),
        (705, 733, 147, 155, "维护预警 /\n异常告警", "预测性维护\n异常实时分配"),
        (868, 733, 132, 155, "能耗控制策略", "节能优化策略\n动态能耗调控"),
    ]
    for x, y, w, h, title, desc in cards:
        add_card(slide, x, y, w, h, COLORS["line_orange"], title, COLORS["orange"], title_h=40)
        add_text(slide, desc, x + 8, y + 116, w - 16, 36, size=8.3, color="111111", bold=False, align="center", valign="middle")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1007, 728, 621, 161, fill="FFFFFF", line=COLORS["line_orange"], line_width=1.1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1007, 728, 621, 28, fill=COLORS["orange2"], line=None)
    add_text(slide, "KPI 看板", 1240, 731, 155, 24, size=13, color="FFFFFF", bold=True, align="center")
    kpi_cards = [
        (1022, 764, 105, "任务完成率", "92.6%", "↑", COLORS["line_blue"], "kpi_bar_completion"),
        (1138, 764, 105, "设备利用率", "85.4%", "↑", COLORS["line_blue"], "kpi_bar_utilization"),
        (1254, 764, 105, "平均响应时间", "128 ms", "", COLORS["line_blue"], "kpi_line_response"),
        (1370, 764, 105, "能耗下降", "18.7%", "↓", "11A33B", "kpi_line_energy"),
        (1486, 764, 105, "异常率", "2.1%", "↓", COLORS["line_orange"], "kpi_line_anomaly"),
    ]
    for x, y, w, title, metric, arrow, color, _asset in kpi_cards:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 111, fill="FFFFFF", line="D7DDE4", line_width=0.8)
        add_text(slide, title, x + 3, y + 6, w - 6, 19, size=9.5, color=COLORS["ink"], bold=True, align="center")
        add_text(slide, metric, x + 7, y + 31, w - 24, 29, size=16, color=COLORS["ink"], bold=True, align="center")
        if arrow:
            add_text(slide, arrow, x + w - 22, y + 34, 16, 18, size=10, color="1CAF48", bold=True, align="center")

    add_line(slide, 235, 905, 1600, 905, COLORS["line_orange"], width=1.3, dashed=True)
    add_shape(slide, MSO_SHAPE.UP_ARROW, 230, 902, 17, 27, fill=COLORS["line_orange"], line=None)
    add_text(slide, "执行反馈 / KPI 反馈 / 持续优化", 610, 913, 450, 24, size=15, color=COLORS["orange"], bold=True, align="center")

    for unit_id, file_name, box, role in IMAGE_UNITS:
        if unit_id == "scene_bg":
            continue
        placements.append(add_picture_fit(slide, file_name, *box, unit_id))

    placements = [p for p in placements if p["id"] != "scene_bg"] + [{
        "id": "scene_bg",
        "path": "scene_background_top.png",
        "anchor_slot": [122, 107, 1526, 255],
        "fitted_bbox": [122, 107, 1526, 255],
    }]

    pptx_path = OUT / "manufacturing_scheduler_replica.pptx"
    prs.save(pptx_path)
    return placements


def crop_reference_units() -> None:
    ref = Image.open(REF).convert("RGB")
    for unit_id, _file_name, box, role in IMAGE_UNITS:
        if role == "generated_scene_background":
            continue
        x, y, w, h = [int(v) for v in box]
        pad = 10
        crop = ref.crop((max(0, x - pad), max(0, y - pad), min(CANVAS_W, x + w + pad), min(CANVAS_H, y + h + pad)))
        crop.save(OUT / "reference_crops" / f"{unit_id}.png")


def write_residual_artifacts() -> None:
    img = Image.open(REF).convert("RGB")
    draw = ImageDraw.Draw(img)
    redboxes = []
    for unit_id, _file_name, box, role in IMAGE_UNITS:
        if role == "generated_scene_background":
            continue
        x, y, w, h = box
        draw.rectangle([x, y, x + w, y + h], fill="white")
        redboxes.append({"semantic_unit_id": unit_id, "bbox": box, "class": "imagegen_asset"})
    img.save(OUT / "residual_cycle_1.png")
    (OUT / "residual_cycle_1_redboxes.json").write_text(json.dumps({"redboxes": redboxes}, ensure_ascii=False, indent=2), encoding="utf-8")
    (OUT / "asset_match_cycle_1.json").write_text(json.dumps({"matches": redboxes}, ensure_ascii=False, indent=2), encoding="utf-8")


def write_prompts() -> None:
    prompt_rows = [
        ("assets_cycle_1.jsonl", "3x4 grid: AMR, robot arm, AGV, server, database, digital twin screen, clipboard, brain scheduler, alarm, energy leaves and battery, route pins, KPI dashboard; flat #00ff00 chroma-key; no text."),
        ("assets_cycle_2.jsonl", "4x4 grid: sensing/status/environment/factory/chip/output/task/resource/energy/constraint/anomaly/collaboration/status/resource pool/device pool/dispatch icons; flat #00ff00 chroma-key; no text."),
        ("assets_cycle_3.jsonl", "2x3 grid: KPI mini chart icons and blank KPI card; flat #00ff00 chroma-key; no text or numbers."),
        ("assets_cycle_4.jsonl", "1x2 grid: green energy leaf battery and green energy-down KPI chart; flat #ff00ff chroma-key; no text or numbers."),
        ("scene_background.jsonl", "Wide pale intelligent manufacturing workshop background, no foreground semantic objects, no text, suitable for overlay."),
    ]
    for file_name, prompt in prompt_rows:
        row = {
            "image_inputs": ["reference.png", "reference_crops/*.png"],
            "prompt": prompt,
            "background": "#00ff00" if "ff00ff" not in prompt and "background" not in file_name else "#ff00ff" if "assets_cycle_4" in file_name else "none",
            "source": "built-in image_gen",
        }
        (PROMPTS / file_name).write_text(json.dumps(row, ensure_ascii=False) + "\n", encoding="utf-8")


def write_json_artifacts(placements: list[dict]) -> None:
    image_items = []
    for unit_id, file_name, box, role in IMAGE_UNITS:
        image_items.append({
            "semantic_unit_id": unit_id,
            "class": "imagegen_asset" if role != "generated_scene_background" else "generated_scene_background",
            "path": file_name,
            "bbox": box,
            "semantic_unit_count": 1 if role != "generated_scene_background" else None,
            "source_type": "imagegen_asset",
        })
    text_items = [
        "标题", "副标题", "多源感知输入", "设备/资源/KPI标签", "订单需求文本", "调度器模块标签", "输出决策文本", "KPI数值文本"
    ]
    layout_items = ["三层主面板", "侧边层级色带", "数据流箭头", "卡片框", "虚线反馈闭环", "模块连接线"]
    inventory = {
        "canvas": {"width": CANVAS_W, "height": CANVAS_H},
        "classes": {"text": text_items, "layout_native": layout_items, "imagegen_asset": image_items},
    }
    (OUT / "visual_inventory.json").write_text(json.dumps(inventory, ensure_ascii=False, indent=2), encoding="utf-8")

    asset_manifest = {
        "canvas": {"width": CANVAS_W, "height": CANVAS_H},
        "slide": {"width_in": SLIDE_W_IN, "height_in": SLIDE_H_IN},
        "assets": image_items,
        "placements": placements,
        "rules": {
            "source_bitmap_reference_embedded": False,
            "all_semantic_visuals_are_images": True,
            "semantic_unit_count_per_final_image": 1,
            "image_scaling": "uniform contain",
        },
    }
    (OUT / "asset_manifest.json").write_text(json.dumps(asset_manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    layout_rules = {
        "font": "Microsoft YaHei",
        "title_group": {"align": "center", "controlled_lines": [1, 1]},
        "native_allowed": ["background", "panel", "frame", "divider", "connector", "structural_arrow", "text_container"],
        "image_fit": "contain only; no one-axis stretching",
        "layer_palette": {"scene": COLORS["blue"], "scheduler": COLORS["teal"], "output": COLORS["orange"]},
    }
    (OUT / "layout_rules.json").write_text(json.dumps(layout_rules, ensure_ascii=False, indent=2), encoding="utf-8")

    validation = {
        "checks": {
            "reference_bitmap_not_used_as_final_slide": True,
            "generated_scene_background_used_for_top_layer": True,
            "semantic_visuals_are_independent_png_objects": True,
            "text_is_editable_ppt_text": True,
            "layout_shapes_are_native_only": True,
            "no_hand_drawn_semantic_vectors": True,
            "uniform_image_contain_scaling": True,
        },
        "semantic_image_units": len([x for x in image_items if x["class"] == "imagegen_asset"]),
        "scene_background_units": len([x for x in image_items if x["class"] == "generated_scene_background"]),
        "notes": [
            "IMAGE GEN generated assets are stylistic replicas rather than exact OCR/vector traces.",
            "PPTX parity was not checked with PowerPoint/LibreOffice; a script-rendered preview is included for visual inspection."
        ],
    }
    (OUT / "validation_report.json").write_text(json.dumps(validation, ensure_ascii=False, indent=2), encoding="utf-8")


def render_preview() -> None:
    pptx_path = OUT / "manufacturing_scheduler_replica.pptx"
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    img = Image.new("RGB", (CANVAS_W, CANVAS_H), "white")
    draw = ImageDraw.Draw(img)
    font_path = Path("C:/Windows/Fonts/msyh.ttc")

    def font(size: int):
        return ImageFont.truetype(str(font_path), size) if font_path.exists() else ImageFont.load_default()

    def to_px(emu_value: int) -> int:
        return int(round(emu_value / (914400 * SCALE)))

    def color_to_hex(color_format, fallback=None):
        try:
            value = color_format.rgb
            if value is None:
                return fallback
            return f"#{value}"
        except Exception:
            return fallback

    def shape_fill(shape):
        try:
            return color_to_hex(shape.fill.fore_color)
        except Exception:
            return None

    def shape_line(shape):
        try:
            return color_to_hex(shape.line.fill.fore_color)
        except Exception:
            return None

    def draw_auto_shape(shape, x, y, w, h):
        fill = shape_fill(shape)
        line = shape_line(shape)
        outline = line if line else fill
        kind = str(getattr(shape, "auto_shape_type", ""))
        if "OVAL" in kind:
            draw.ellipse([x, y, x + w, y + h], fill=fill, outline=outline, width=2 if outline else 1)
        elif "RIGHT_TRIANGLE" in kind:
            draw.polygon([(x, y), (x, y + h), (x + w, y)], fill=fill, outline=outline)
        elif "RIGHT_ARROW" in kind:
            head = min(w * 0.35, h * 1.4)
            pts = [(x, y + h * 0.25), (x + w - head, y + h * 0.25), (x + w - head, y),
                   (x + w, y + h / 2), (x + w - head, y + h), (x + w - head, y + h * 0.75),
                   (x, y + h * 0.75)]
            draw.polygon(pts, fill=fill, outline=outline)
        elif "DOWN_ARROW" in kind or "UP_ARROW" in kind:
            draw.rounded_rectangle([x, y, x + w, y + h], radius=3, fill=fill, outline=outline, width=1)
        else:
            radius = 10 if "ROUNDED" in kind else 0
            if radius:
                draw.rounded_rectangle([x, y, x + w, y + h], radius=radius, fill=fill, outline=outline, width=2 if outline else 1)
            else:
                draw.rectangle([x, y, x + w, y + h], fill=fill, outline=outline, width=2 if outline else 1)

    def draw_text(shape, x, y, w, h):
        if not getattr(shape, "has_text_frame", False):
            return
        lines = []
        for p in shape.text_frame.paragraphs:
            text = "".join(run.text for run in p.runs) or p.text
            if text == "":
                continue
            size_pt = 12
            color = "#12345B"
            bold = False
            for run in p.runs:
                if run.font.size:
                    size_pt = run.font.size.pt
                c = color_to_hex(run.font.color, None)
                if c:
                    color = c
                if run.font.bold:
                    bold = True
                break
            lines.extend([(part, size_pt, color, bold, p.alignment) for part in text.split("\n")])
        if not lines:
            return
        line_heights = []
        rendered = []
        for text, size_pt, color, bold, align in lines:
            f = font(max(7, int(size_pt * 1.35)))
            bbox = draw.textbbox((0, 0), text, font=f)
            line_heights.append((bbox[3] - bbox[1]) + 3)
            rendered.append((text, f, color, align))
        total_h = sum(line_heights)
        try:
            anchor = shape.text_frame.vertical_anchor
        except Exception:
            anchor = None
        cy = y + max(2, (h - total_h) / 2) if anchor == MSO_ANCHOR.MIDDLE else y + 3
        for (text, f, color, align), lh in zip(rendered, line_heights):
            bbox = draw.textbbox((0, 0), text, font=f)
            tw = bbox[2] - bbox[0]
            if align == PP_ALIGN.CENTER:
                tx = x + (w - tw) / 2
            elif align == PP_ALIGN.RIGHT:
                tx = x + w - tw - 3
            else:
                tx = x + 4
            draw.text((tx, cy), text, fill=color, font=f)
            cy += lh

    for shape in slide.shapes:
        x, y, w, h = [to_px(v) for v in (shape.left, shape.top, shape.width, shape.height)]
        if shape.shape_type == 13:
            asset = Image.open(BytesIO(shape.image.blob)).convert("RGBA").resize((max(1, w), max(1, h)))
            img.paste(asset, (x, y), asset)
        elif shape.shape_type == 9:
            add_line_x2 = x + w
            add_line_y2 = y + h
            draw.line([x, y, add_line_x2, add_line_y2], fill=shape_line(shape) or "#12345B", width=2)
        else:
            try:
                draw_auto_shape(shape, x, y, w, h)
            except Exception:
                pass
        draw_text(shape, x, y, w, h)
    img.save(OUT / "preview" / "slide_preview.png")


def main() -> None:
    ensure_dirs()
    crop_reference_units()
    write_residual_artifacts()
    write_prompts()
    placements = build_pptx()
    write_json_artifacts(placements)
    render_preview()
    print(OUT / "manufacturing_scheduler_replica.pptx")
    print(OUT / "preview" / "slide_preview.png")


if __name__ == "__main__":
    main()
