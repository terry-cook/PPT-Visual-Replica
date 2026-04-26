from __future__ import annotations

import json
import math
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
WORKSPACE = SCRIPT_DIR.parent
REPO_ROOT = WORKSPACE.parents[1]
OUT = REPO_ROOT / "local-runs" / "satellite-network-rebuild"
GRID_SOURCE = WORKSPACE / "imagegen" / "generated" / "imagegen_asset_grid_cycle_1.png"
SPACE_SOURCE = WORKSPACE / "imagegen" / "generated" / "imagegen_space_background.png"
REF_SOURCE = WORKSPACE / "reference" / "reference.png"

SLIDE_W_PX = 1680
SLIDE_H_PX = 945
PX_PER_IN = 105

NAVY = "0B2B78"
BLUE = "105BC9"
MID_BLUE = "1E73D8"
LIGHT_BLUE = "EAF4FF"
PALE_BLUE = "F4FAFF"
ORANGE = "EA6B1D"
GREEN = "3F9D43"
TEAL = "0E8D98"
INK = "082A6D"
GRAY = "6E7C95"


ASSET_NAMES = [
    "satellite_node",
    "ground_station",
    "database_stack",
    "bar_chart",
    "line_pie_chart",
    "network_graph",
    "shield_check",
    "robot_agent",
    "server_rack",
    "document_check",
    "monitor_dashboard",
    "target",
    "circular_arrows",
    "wireless_status",
    "collaboration_group",
    "brain_head",
    "server_gear",
    "task_cube",
]


def ensure_dirs() -> None:
    for rel in [
        "",
        "reference_crops",
        "generated",
        "assets",
        "prompts",
        "previews",
    ]:
        (OUT / rel).mkdir(parents=True, exist_ok=True)


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def px(value: float):
    return Inches(value / PX_PER_IN)


def cut_chroma_grid() -> dict[str, Path]:
    dst_grid = OUT / "generated" / "imagegen_asset_grid_cycle_1.png"
    dst_space = OUT / "generated" / "imagegen_space_background.png"
    shutil.copy2(GRID_SOURCE, dst_grid)
    shutil.copy2(SPACE_SOURCE, dst_space)

    img = Image.open(dst_grid).convert("RGBA")
    w, h = img.size
    cols, rows = 6, 3
    cell_w, cell_h = w / cols, h / rows
    assets: dict[str, Path] = {}

    for idx, name in enumerate(ASSET_NAMES):
        c = idx % cols
        r = idx // cols
        left = int(c * cell_w)
        upper = int(r * cell_h)
        right = int((c + 1) * cell_w)
        lower = int((r + 1) * cell_h)
        cell = img.crop((left, upper, right, lower)).convert("RGBA")
        data = cell.load()
        for y in range(cell.height):
            for x in range(cell.width):
                rr, gg, bb, aa = data[x, y]
                if gg > 150 and rr < 120 and bb < 120:
                    data[x, y] = (0, 0, 0, 0)

        alpha = cell.getchannel("A")
        bbox = alpha.getbbox()
        if bbox:
            pad = 8
            bbox = (
                max(0, bbox[0] - pad),
                max(0, bbox[1] - pad),
                min(cell.width, bbox[2] + pad),
                min(cell.height, bbox[3] + pad),
            )
            cell = cell.crop(bbox)

        out = OUT / "assets" / f"{name}.png"
        cell.save(out)
        assets[name] = out

    bg = Image.open(dst_space).convert("RGB")
    # Keep the horizon visible in the short banner crop.
    crop = bg.crop((0, int(bg.height * 0.42), bg.width, int(bg.height * 0.88)))
    crop = ImageOps.fit(crop, (1320, 175), method=Image.Resampling.LANCZOS)
    strip = OUT / "assets" / "space_earth_background.png"
    crop.save(strip)
    assets["space_earth_background"] = strip

    make_tinted(assets["bar_chart"], OUT / "assets" / "bar_chart_orange.png", ORANGE)
    make_tinted(assets["bar_chart"], OUT / "assets" / "bar_chart_green.png", GREEN)
    make_tinted(assets["network_graph"], OUT / "assets" / "network_graph_orange.png", ORANGE)
    make_tinted(assets["server_gear"], OUT / "assets" / "server_gear_teal.png", TEAL)
    assets["bar_chart_orange"] = OUT / "assets" / "bar_chart_orange.png"
    assets["bar_chart_green"] = OUT / "assets" / "bar_chart_green.png"
    assets["network_graph_orange"] = OUT / "assets" / "network_graph_orange.png"
    assets["server_gear_teal"] = OUT / "assets" / "server_gear_teal.png"
    return assets


def make_tinted(src: Path, dst: Path, hex_color: str) -> None:
    base = Image.open(src).convert("RGBA")
    target = tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    pix = base.load()
    for y in range(base.height):
        for x in range(base.width):
            r, g, b, a = pix[x, y]
            if a == 0:
                continue
            lum = int(0.299 * r + 0.587 * g + 0.114 * b)
            mix = lum / 255
            nr = int(target[0] * (0.45 + 0.45 * mix) + 255 * 0.10 * mix)
            ng = int(target[1] * (0.45 + 0.45 * mix) + 255 * 0.10 * mix)
            nb = int(target[2] * (0.45 + 0.45 * mix) + 255 * 0.10 * mix)
            pix[x, y] = (min(255, nr), min(255, ng), min(255, nb), a)
    base.save(dst)


def add_box(slide, x, y, w, h, fill="FFFFFF", line=MID_BLUE, radius=True, width=1.2):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(width)
    return shp


def add_text(slide, text, x, y, w, h, size=18, color=NAVY, bold=False, align="center", valign="mid"):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = px(4)
    tf.margin_right = px(4)
    tf.margin_top = px(1)
    tf.margin_bottom = px(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if valign == "mid" else MSO_ANCHOR.TOP
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}[align]
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return box


def add_pic(slide, path: Path, x, y, w, h, placements, semantic_id: str, contain=True):
    img = Image.open(path)
    iw, ih = img.size
    bx, by, bw, bh = x, y, w, h
    if contain:
        scale = min(w / iw, h / ih)
        fw = iw * scale
        fh = ih * scale
        bx = x + (w - fw) / 2
        by = y + (h - fh) / 2
        bw, bh = fw, fh
    slide.shapes.add_picture(str(path), px(bx), px(by), width=px(bw), height=px(bh))
    placements.append(
        {
            "semantic_unit_id": semantic_id,
            "asset": path.name,
            "bbox_px": [round(bx, 1), round(by, 1), round(bw, 1), round(bh, 1)],
            "semantic_unit_count": 1,
            "fit": "contain",
        }
    )


def add_arrow(slide, x, y, w, h, color=MID_BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.color.rgb = rgb(color)
    return shp


def add_down_arrow(slide, x, y, w, h, color=MID_BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.color.rgb = rgb(color)
    return shp


def add_line(slide, x1, y1, x2, y2, color=MID_BLUE, width=1.2, dash=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = 4
    return line


def build_deck(assets: dict[str, Path], placements: list[dict]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("F7FBFF")

    add_text(slide, "面向异构卫星网络的多智能体协同卸载与资源调度总体架构", 180, 10, 1320, 55, 29, NAVY, True)
    add_text(slide, "融合业务态势感知、低交互协同卸载与仿真验证的研究框架", 455, 70, 770, 36, 17, NAVY, True)

    # Top background and application scenario strip.
    add_box(slide, 18, 110, 166, 158, fill="0758BA", line=MID_BLUE, width=1.0)
    add_text(slide, "研究背景\n与应用场景", 35, 150, 128, 70, 20, "FFFFFF", True)
    add_box(slide, 184, 110, 1305, 158, fill="FFFFFF", line=MID_BLUE, width=1.0)
    add_pic(slide, assets["space_earth_background"], 185, 112, 1300, 154, placements, "scene_top_space_background", contain=False)

    top_cards = [
        ("satellite_node", "异构卫星网络规模扩大，\n星间协作关系复杂", 270, 128, 300, 68),
        ("database_stack", "多源状态信息分散，\n业务态势动态演化", 670, 128, 280, 68),
        ("bar_chart", "任务负载波动明显，\n资源调度稳定性要求高", 1060, 128, 330, 68),
    ]
    for i, (icon, text, x, y, w, h) in enumerate(top_cards):
        add_box(slide, x, y, w, h, fill="F6FBFF", line="7EA8EF", width=1.0)
        add_pic(slide, assets[icon], x + 18, y + 13, 44, 42, placements, f"top_card_icon_{i+1}")
        add_text(slide, text, x + 74, y + 11, w - 88, h - 18, 14.5, NAVY, True, "left")

    for i, (x, y, sz) in enumerate([(230, 220, 34), (570, 224, 34), (845, 218, 28), (1010, 198, 28), (1115, 220, 32), (1390, 220, 44)]):
        add_pic(slide, assets["ground_station"] if i in (0, 5) else assets["satellite_node"], x, y, sz, sz, placements, f"top_network_node_{i+1}")
    for i, (x, y) in enumerate([(385, 214), (720, 212), (1200, 214)]):
        add_pic(slide, assets["task_cube"], x, y, 20, 20, placements, f"top_cube_{i+1}")
    add_line(slide, 250, 232, 570, 224, "FFFFFF", 1.3, True)
    add_line(slide, 570, 224, 850, 218, "FFFFFF", 1.3, True)
    add_line(slide, 850, 218, 1118, 224, "FFFFFF", 1.3, True)
    add_line(slide, 1118, 224, 1395, 224, "FFFFFF", 1.3, True)
    add_line(slide, 275, 239, 520, 218, ORANGE, 1.3, True)
    add_line(slide, 1135, 232, 1375, 244, ORANGE, 1.3, True)

    add_box(slide, 1490, 110, 172, 158, fill="F8FBFF", line="7EA8EF", width=1.0)
    legend_items = [
        ("satellite_node", "卫星节点"),
        ("line", "星间链路"),
        ("task_cube", "任务数据流"),
        ("orange", "业务态势流"),
        ("ground_station", "地面站"),
    ]
    y0 = 128
    for i, (kind, label) in enumerate(legend_items):
        yy = y0 + i * 26
        if kind == "line":
            add_line(slide, 1518, yy + 10, 1562, yy + 10, NAVY, 1.3, True)
        elif kind == "orange":
            add_line(slide, 1518, yy + 10, 1562, yy + 10, ORANGE, 1.3, True)
        else:
            add_pic(slide, assets[kind], 1518, yy, 26, 22, placements, f"legend_{kind}")
        add_text(slide, label, 1565, yy - 2, 82, 25, 9.5, NAVY, False, "left")

    build_content_rows(slide, assets, placements)

    out_pptx = OUT / "satellite_network_framework_replica.pptx"
    prs.save(out_pptx)
    return out_pptx


def build_content_rows(slide, assets: dict[str, Path], placements: list[dict]) -> None:
    # Row 1
    add_box(slide, 18, 284, 234, 156, fill="0758BA", line=MID_BLUE)
    add_text(slide, "研究内容 1：\n基于业务态势时空扩展图\n的关联机理分析", 32, 300, 210, 78, 13.7, "FFFFFF", True, "left")
    add_pic(slide, assets["network_graph"], 82, 382, 86, 44, placements, "row1_left_network")
    add_box(slide, 250, 284, 1288, 156, fill="FFFFFF", line=MID_BLUE)

    row1 = [
        ("多源状态信息融合", "satellite_node", 310, 298, 205, 106),
        ("业务态势感知", "line_pie_chart", 580, 298, 170, 106),
        ("动态协作计算\n关系建模", "network_graph", 815, 298, 180, 106),
        ("时空扩展图构建", "network_graph", 1060, 298, 184, 106),
        ("静态图扩展的\n可靠转换路径", "shield_check", 1310, 298, 170, 106),
    ]
    for i, (title, icon, x, y, w, h) in enumerate(row1):
        add_box(slide, x, y, w, h, fill=PALE_BLUE, line="7EA8EF")
        add_text(slide, title, x + 10, y + 8, w - 20, 32, 13.5, NAVY, True)
        add_pic(slide, assets[icon], x + 42, y + 48, w - 84, 46, placements, f"row1_module_{i+1}_{icon}")
        if i < len(row1) - 1:
            add_arrow(slide, x + w + 15, y + 44, 38, 22)
    add_text(slide, "目标：揭示业务态势与协作关系的关联机理，为卸载决策提供结构化表征", 520, 410, 760, 22, 11.2, NAVY, True)
    add_down_arrow(slide, 805, 262, 42, 36)
    add_down_arrow(slide, 805, 432, 42, 34)
    add_text(slide, "表征支撑决策", 872, 440, 190, 18, 9.6, NAVY, True)

    # Row 2
    add_box(slide, 18, 455, 234, 170, fill="0758BA", line=MID_BLUE)
    add_text(slide, "研究内容 2：\n基于多智能体强化学习的\n低交互协同卸载技术", 32, 472, 210, 88, 13.5, "FFFFFF", True, "left")
    add_pic(slide, assets["collaboration_group"], 84, 564, 84, 42, placements, "row2_left_group")
    add_box(slide, 250, 455, 1288, 170, fill="FFFFFF", line=MID_BLUE)

    add_box(slide, 340, 474, 200, 122, fill="F4F8FF", line="7EA8EF")
    inputs = [("document_check", "任务需求"), ("server_rack", "资源状态"), ("wireless_status", "网络状态")]
    for i, (icon, label) in enumerate(inputs):
        yy = 488 + i * 38
        add_pic(slide, assets[icon], 365, yy, 28, 28, placements, f"row2_input_{i+1}")
        add_text(slide, label, 414, yy + 2, 92, 24, 12.8, NAVY, True, "left")
    add_arrow(slide, 560, 520, 58, 28)

    add_box(slide, 640, 474, 480, 122, fill="F7FBFF", line="7EA8EF")
    add_text(slide, "多智能体强化学习协同决策", 720, 480, 320, 26, 15, NAVY, True)
    agent_x = [700, 815, 940, 1050]
    for i, ax in enumerate(agent_x):
        add_pic(slide, assets["robot_agent"], ax, 510, 42, 38, placements, f"row2_agent_{i+1}")
        if i < len(agent_x) - 1:
            add_line(slide, ax + 44, 528, agent_x[i + 1] - 5, 528, NAVY, 1.1, True)
    labels = [("分布式决策", 655, ORANGE), ("低交互开销", 833, ORANGE), ("高效协作", 1000, ORANGE)]
    for label, x, color in labels:
        add_box(slide, x, 556, 132, 30, fill="FFF7F2", line=color, width=1.0)
        add_text(slide, label, x, 558, 132, 24, 12.5, color, True)
    add_arrow(slide, 1138, 520, 58, 28)

    add_box(slide, 1215, 474, 240, 122, fill="F0FFFB", line="70B9C0")
    add_pic(slide, assets["document_check"], 1240, 490, 36, 34, placements, "row2_output_strategy")
    add_text(slide, "协同卸载策略", 1295, 493, 112, 28, 12.5, NAVY, True, "left")
    add_pic(slide, assets["server_rack"], 1240, 542, 36, 34, placements, "row2_output_resource")
    add_text(slide, "资源分配方案", 1295, 545, 112, 28, 12.5, NAVY, True, "left")
    add_text(slide, "目标：在低交互条件下实现多智能体高效协同卸载", 580, 604, 580, 22, 11.2, NAVY, True)
    add_down_arrow(slide, 805, 617, 42, 34)
    add_text(slide, "策略嵌入仿真", 870, 628, 170, 16, 9.6, NAVY, True)

    # Row 3
    add_box(slide, 18, 640, 234, 160, fill="0758BA", line=MID_BLUE)
    add_text(slide, "研究内容 3：\n多智能体资源协同\n调度仿真验证系统", 32, 656, 210, 82, 13.8, "FFFFFF", True, "left")
    add_pic(slide, assets["monitor_dashboard"], 78, 742, 90, 44, placements, "row3_left_monitor")
    add_box(slide, 250, 640, 1288, 160, fill="FFFFFF", line=MID_BLUE)

    add_box(slide, 275, 650, 360, 138, fill=PALE_BLUE, line="7EA8EF")
    add_text(slide, "异构卫星智能体集群", 330, 654, 250, 28, 15, NAVY, True)
    sat_positions = [(310, 690), (380, 690), (450, 690), (520, 690), (590, 690), (342, 725), (450, 725), (560, 725)]
    for i, (sx, sy) in enumerate(sat_positions):
        add_pic(slide, assets["satellite_node"], sx, sy, 38, 30, placements, f"row3_satellite_{i+1}")
    add_text(slide, "种类不少于3种，数量不少于50个", 360, 762, 240, 20, 11.5, NAVY, True)

    add_box(slide, 690, 650, 335, 138, fill="F9FBFF", line="7EA8EF")
    add_text(slide, "任务负载场景", 770, 654, 180, 28, 15, NAVY, True)
    loads = [
        ("高负载", "bar_chart_orange", ORANGE, 708),
        ("中负载", "bar_chart", NAVY, 815),
        ("低负载", "bar_chart_green", GREEN, 922),
    ]
    for label, icon, color, x in loads:
        add_box(slide, x, 690, 90, 78, fill="FFF9F5" if color == ORANGE else "F8FBFF", line="BFCDEB")
        add_text(slide, label, x, 696, 90, 22, 12.5, color, True)
        add_pic(slide, assets[icon], x + 22, 724, 44, 34, placements, f"row3_load_{label}")
    add_text(slide, "...", 1030, 710, 42, 40, 20, NAVY, True)

    add_box(slide, 1070, 650, 460, 138, fill=PALE_BLUE, line="7EA8EF")
    add_text(slide, "仿真评估指标", 1180, 654, 210, 26, 15, NAVY, True)
    metrics = [("circular_arrows", "任务卸载仿真验证"), ("wireless_status", "协同调度性能评估"), ("shield_check", "性能波动不超过10%")]
    for i, (icon, label) in enumerate(metrics):
        yy = 690 + i * 30
        add_pic(slide, assets[icon], 1092, yy, 28, 26, placements, f"row3_metric_{i+1}")
        add_text(slide, label, 1130, yy + 1, 190, 22, 11.5, NAVY, True, "left")
    add_pic(slide, assets["monitor_dashboard"], 1320, 666, 170, 95, placements, "row3_dashboard_monitor")
    add_text(slide, "目标：验证所提智能卸载方案在不同负载场景下的鲁棒性与稳定性", 555, 796, 690, 24, 12.5, NAVY, True)

    # Feedback loop on the right, structural native arrow.
    arrow = slide.shapes.add_shape(MSO_SHAPE.CURVED_DOWN_ARROW, px(1535), px(345), px(96), px(410))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(MID_BLUE)
    arrow.line.color.rgb = rgb(MID_BLUE)
    add_text(slide, "验证反馈\n优化", 1580, 505, 75, 62, 13, NAVY, True)

    # Bottom output row.
    add_box(slide, 18, 835, 248, 96, fill="0758BA", line=MID_BLUE)
    add_text(slide, "总体产出 /\n预期目标", 42, 852, 142, 52, 20, "FFFFFF", True, "left")
    add_pic(slide, assets["target"], 180, 845, 62, 62, placements, "bottom_target")
    add_box(slide, 266, 835, 1300, 96, fill="FFFFFF", line=MID_BLUE)
    flow = [
        ("brain_head", "态势认知增强", 300, 852, 240, BLUE),
        ("network_graph_orange", "协同卸载优化", 630, 852, 240, ORANGE),
        ("server_gear_teal", "资源调度稳定", 950, 852, 250, TEAL),
        ("circular_arrows", "系统仿真验证闭环", 1275, 852, 250, BLUE),
    ]
    for i, (icon, label, x, y, w, color) in enumerate(flow):
        add_box(slide, x, y, w, 52, fill="F8FBFF", line=color, width=1.2)
        add_pic(slide, assets[icon], x + 18, y + 8, 40, 36, placements, f"bottom_flow_{i+1}_{icon}")
        add_text(slide, label, x + 62, y + 8, w - 72, 34, 14.5, color, True, "left")
        if i < len(flow) - 1:
            add_arrow(slide, x + w + 34, y + 16, 48, 22, "87A9D8")
    add_text(slide, "形成从态势感知、决策优化到系统验证的一体化技术体系", 555, 904, 620, 24, 12.5, NAVY, True)


def create_reference_crops() -> None:
    ref = Image.open(OUT / "reference.png").convert("RGB")
    crops = {
        "satellite_node": [292, 664, 350, 720],
        "ground_station": [215, 215, 260, 255],
        "database_stack": [420, 340, 470, 390],
        "bar_chart": [1015, 722, 1080, 770],
        "line_pie_chart": [605, 335, 720, 390],
        "network_graph": [835, 345, 968, 390],
        "shield_check": [1425, 350, 1470, 395],
        "robot_agent": [710, 505, 750, 552],
        "server_rack": [363, 508, 390, 558],
        "document_check": [1225, 492, 1272, 536],
        "monitor_dashboard": [1320, 655, 1510, 775],
        "target": [182, 842, 244, 905],
        "circular_arrows": [1296, 858, 1338, 898],
        "wireless_status": [365, 555, 395, 588],
        "collaboration_group": [85, 560, 170, 612],
        "brain_head": [316, 858, 355, 896],
        "server_gear": [965, 858, 1010, 896],
        "task_cube": [384, 214, 404, 234],
    }
    redboxes = []
    draw = ImageDraw.Draw(ref)
    for name, box in crops.items():
        crop = ref.crop(tuple(box))
        crop.save(OUT / "reference_crops" / f"{name}.png")
        redboxes.append({"semantic_unit_id": name, "bbox_px": box, "class": "imagegen_asset"})
        draw.rectangle(tuple(box), outline=(255, 0, 0), width=3)
    ref.save(OUT / "residual_cycle_1.png")
    (OUT / "residual_cycle_1_redboxes.json").write_text(
        json.dumps(redboxes, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def write_artifacts(assets: dict[str, Path], placements: list[dict], out_pptx: Path) -> None:
    asset_records = []
    for name, path in sorted(assets.items()):
        role = "generated_scene_background" if name == "space_earth_background" else "imagegen_asset"
        asset_records.append(
            {
                "asset_id": name,
                "role": role,
                "path": str(path.relative_to(OUT)).replace("\\", "/"),
                "source": "imagegen_asset_grid_cycle_1" if role == "imagegen_asset" else "imagegen_space_background",
                "semantic_unit_count": 1 if role == "imagegen_asset" else "background",
            }
        )

    inventory = {
        "reference": "reference.png",
        "slide_size_px": [SLIDE_W_PX, SLIDE_H_PX],
        "classes": {
            "text": "All Chinese titles, labels, subtitles, and objectives are native editable PPT text boxes.",
            "layout_native": "Panels, section boxes, dividers, flow arrows, dashed links, and structural connectors are native PPT objects.",
            "imagegen_asset": "All semantic icons, pictograms, dashboards, charts, satellites, stations, robots, and cubes are imagegen-derived PNG assets.",
        },
        "minimum_semantic_units": asset_records,
        "placements": placements,
    }
    (OUT / "visual_inventory.json").write_text(json.dumps(inventory, ensure_ascii=False, indent=2), encoding="utf-8")
    (OUT / "asset_manifest.json").write_text(
        json.dumps({"assets": asset_records, "placements": placements}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    (OUT / "asset_match_cycle_1.json").write_text(
        json.dumps(
            {
                "imagegen_grid": "generated/imagegen_asset_grid_cycle_1.png",
                "background": "generated/imagegen_space_background.png",
                "cut_rule": "3 rows x 6 columns, equal cell split, chroma-key #00ff00 removed, alpha trim applied",
                "matched_assets": ASSET_NAMES,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    (OUT / "layout_rules.json").write_text(
        json.dumps(
            {
                "coordinate_system": "Reference pixels mapped to 16:9 slide at 105 px per inch.",
                "image_scaling": "uniform contain; no one-axis stretch for semantic assets",
                "native_allowed": ["text boxes", "panels", "dividers", "structural arrows", "flow connectors"],
                "native_forbidden": ["semantic satellite icons", "chart icons", "robot/user icons", "dashboard screens", "network pictograms"],
                "font": "Microsoft YaHei",
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    prompt = {
        "cycle": 1,
        "prompt_summary": "Generated 18 isolated infographic icons on #00ff00 chroma-key background using the reference slide for style and identity.",
        "grid": "3x6, margin 80 px, gap 48 px, no text",
        "objects_in_order": ASSET_NAMES,
    }
    (OUT / "prompts" / "assets_cycle_1.jsonl").write_text(
        json.dumps(prompt, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    (OUT / "validation_report.json").write_text(
        json.dumps(
            {
                "pptx": str(out_pptx.relative_to(OUT)).replace("\\", "/"),
                "checks": [
                    "Reference bitmap is not inserted into the final PPTX.",
                    "Semantic non-text visuals are inserted as individual transparent PNG picture objects.",
                    "Each placement in asset_manifest has semantic_unit_count = 1 except the approved generated scene background.",
                    "Text remains native editable PPT text.",
                    "Structural arrows, panels, connectors, and dividers are native PPT objects.",
                    "Generated icons were cut from an imagegen grid and locally chroma-keyed to transparent PNG.",
                ],
                "known_limits": [
                    "PowerPoint round-trip render was not performed by this script.",
                    "The generated icons intentionally match semantic role and palette, not pixel-identical source glyphs.",
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )


def font(size: int, bold=False):
    font_paths = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\arial.ttf"),
    ]
    for p in font_paths:
        if p.exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def render_preview(assets: dict[str, Path]) -> Path:
    img = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), "#F7FBFF")
    draw = ImageDraw.Draw(img)

    def rect(x, y, w, h, fill, outline=MID_BLUE, width=2):
        draw.rounded_rectangle([x, y, x + w, y + h], radius=10, fill=f"#{fill}", outline=f"#{outline}", width=width)

    def text(txt, x, y, w, h, size, fill=NAVY, bold=False, align="center"):
        f = font(size, bold)
        lines = txt.split("\n")
        line_h = size + 6
        total_h = line_h * len(lines)
        cy = y + (h - total_h) / 2
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=f)
            tw = bbox[2] - bbox[0]
            tx = x + (w - tw) / 2 if align == "center" else x
            draw.text((tx, cy), line, font=f, fill=f"#{fill}")
            cy += line_h

    def paste_asset(name, x, y, w, h):
        a = Image.open(assets[name]).convert("RGBA")
        a.thumbnail((int(w), int(h)), Image.Resampling.LANCZOS)
        img.paste(a, (int(x + (w - a.width) / 2), int(y + (h - a.height) / 2)), a)

    text("面向异构卫星网络的多智能体协同卸载与资源调度总体架构", 180, 10, 1320, 55, 31, NAVY, True)
    text("融合业务态势感知、低交互协同卸载与仿真验证的研究框架", 455, 70, 770, 36, 19, NAVY, True)
    rect(18, 110, 166, 158, "0758BA")
    text("研究背景\n与应用场景", 35, 150, 128, 70, 24, "FFFFFF", True)
    rect(184, 110, 1305, 158, "FFFFFF")
    bg = Image.open(assets["space_earth_background"]).convert("RGB").resize((1300, 154), Image.Resampling.LANCZOS)
    img.paste(bg, (185, 112))
    for icon, label, x, y, w, h in [
        ("satellite_node", "异构卫星网络规模扩大，\n星间协作关系复杂", 270, 128, 300, 68),
        ("database_stack", "多源状态信息分散，\n业务态势动态演化", 670, 128, 280, 68),
        ("bar_chart", "任务负载波动明显，\n资源调度稳定性要求高", 1060, 128, 330, 68),
    ]:
        rect(x, y, w, h, "F6FBFF", "7EA8EF", 2)
        paste_asset(icon, x + 18, y + 13, 44, 42)
        text(label, x + 78, y + 11, w - 92, h - 18, 16, NAVY, True, "left")

    row_specs = [
        ("研究内容 1：\n基于业务态势时空扩展图\n的关联机理分析", 284, 156, "network_graph"),
        ("研究内容 2：\n基于多智能体强化学习的\n低交互协同卸载技术", 455, 170, "collaboration_group"),
        ("研究内容 3：\n多智能体资源协同\n调度仿真验证系统", 640, 160, "monitor_dashboard"),
        ("总体产出 /\n预期目标", 835, 96, "target"),
    ]
    for label, y, h, icon in row_specs:
        rect(18, y, 234 if y < 835 else 248, h, "0758BA")
        text(label, 34, y + 14, 190, min(88, h - 16), 18 if y < 835 else 22, "FFFFFF", True, "left")
        paste_asset(icon, 78 if y != 835 else 180, y + h - 58 if y != 835 else 845, 90 if y != 835 else 62, 48 if y != 835 else 62)
        rect(250 if y < 835 else 266, y, 1288 if y < 835 else 1300, h, "FFFFFF")

    # Lightweight preview of repeated semantic visuals.
    for name, x, y, w, h in [
        ("database_stack", 360, 340, 52, 45),
        ("line_pie_chart", 625, 340, 80, 48),
        ("network_graph", 850, 340, 115, 52),
        ("shield_check", 1370, 340, 70, 54),
        ("robot_agent", 700, 510, 42, 38),
        ("robot_agent", 815, 510, 42, 38),
        ("robot_agent", 940, 510, 42, 38),
        ("robot_agent", 1050, 510, 42, 38),
        ("monitor_dashboard", 1320, 666, 170, 95),
        ("brain_head", 318, 858, 40, 36),
        ("network_graph_orange", 650, 858, 42, 36),
        ("server_gear_teal", 970, 858, 42, 36),
        ("circular_arrows", 1296, 858, 42, 36),
    ]:
        paste_asset(name, x, y, w, h)

    preview = OUT / "previews" / "layout_preview.png"
    img.save(preview)
    return preview


def main() -> None:
    ensure_dirs()
    shutil.copy2(REF_SOURCE, OUT / "reference.png")
    shutil.copy2(REF_SOURCE, OUT / "residual_cycle_0.png")
    assets = cut_chroma_grid()
    create_reference_crops()
    placements: list[dict] = []
    out_pptx = build_deck(assets, placements)
    preview = render_preview(assets)
    write_artifacts(assets, placements, out_pptx)
    print(json.dumps({"pptx": str(out_pptx), "preview": str(preview)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
