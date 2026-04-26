#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ALLOWED_NATIVE = {"background", "panel", "frame", "divider", "connector", "structural_arrow", "text_container", "layout_only"}


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def emu(px: float, scale: float) -> int:
    return int(Inches(px * scale))


def require(condition: bool, message: str) -> None:
    if not condition:
        raise ValueError(message)


def image_box(item: dict[str, Any], asset_dir: Path) -> tuple[float, float, float, float]:
    slot = item.get("anchor_slot") or [item.get("x"), item.get("y"), item.get("w"), item.get("h")]
    x, y, w, h = [float(v) for v in slot]
    path = asset_dir / str(item["path"])
    with Image.open(path) as img:
        iw, ih = img.size
    fit = min(w / iw, h / ih)
    fw, fh = iw * fit, ih * fit
    return x + (w - fw) / 2, y + (h - fh) / 2, fw, fh


def add_text(slide, item: dict[str, Any], scale: float) -> None:
    shape = slide.shapes.add_textbox(emu(item["x"], scale), emu(item["y"], scale), emu(item["w"], scale), emu(item["h"], scale))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(item.get("valign", "top"), MSO_ANCHOR.TOP)
    for name in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, name, Pt(float(item.get(name, 1))))
    for idx, line in enumerate(str(item.get("text", "")).splitlines() or [""]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(item.get("align", "left"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.name = item.get("font", "Microsoft YaHei")
        run.font.size = Pt(float(item.get("font_size", 12)))
        run.font.bold = bool(item.get("bold", False))
        run.font.color.rgb = rgb(item.get("color", "#111111"))


def add_shape(slide, item: dict[str, Any], scale: float) -> None:
    purpose = str(item.get("purpose", ""))
    require(purpose in ALLOWED_NATIVE, f"{item.get('id')}: native shape purpose is not layout-only: {purpose}")
    shape_type = {"rect": MSO_SHAPE.RECTANGLE, "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE, "right_arrow": MSO_SHAPE.RIGHT_ARROW, "down_arrow": MSO_SHAPE.DOWN_ARROW}.get(item["type"])
    require(shape_type is not None, f"{item.get('id')}: unsupported native shape type")
    shape = slide.shapes.add_shape(shape_type, emu(item["x"], scale), emu(item["y"], scale), emu(item["w"], scale), emu(item["h"], scale))
    if item.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(item["fill"])
    else:
        shape.fill.background()
    if item.get("line"):
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = rgb(item["line"])
        shape.line.width = Pt(float(item.get("line_width", 1)))
    else:
        shape.line.fill.background()


def add_line(slide, item: dict[str, Any], scale: float) -> None:
    require(str(item.get("purpose")) in ALLOWED_NATIVE, f"{item.get('id')}: line purpose must be layout-only")
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(item["x1"], scale), emu(item["y1"], scale), emu(item["x2"], scale), emu(item["y2"], scale))
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = rgb(item.get("color", "#000000"))
    shape.line.width = Pt(float(item.get("width", 1)))


def build(manifest_path: Path, out: Path, asset_dir: Path) -> None:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8-sig"))
    canvas = manifest.get("canvas", {"width": 1600, "height": 900})
    slide_size = manifest.get("slide", {"width_in": 16, "height_in": 9})
    scale = float(slide_size["width_in"]) / float(canvas["width"])
    prs = Presentation()
    prs.slide_width = Inches(float(slide_size["width_in"]))
    prs.slide_height = Inches(float(slide_size["height_in"]))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for item in manifest.get("elements", []):
        typ = item.get("type")
        if typ == "text":
            add_text(slide, item, scale)
        elif typ == "image":
            require(item.get("source_type") in {"imagegen_asset", "openai_image", "gemini_image", "user_asset"}, f"{item.get('id')}: invalid image source")
            is_scene = item.get("asset_role") == "generated_scene_background" and item.get("allows_text_overlay") is True
            if not is_scene:
                require(item.get("semantic_unit_id"), f"{item.get('id')}: image element must declare semantic_unit_id")
                require(int(item.get("semantic_unit_count", 1)) == 1, f"{item.get('id')}: final image asset must contain exactly one minimum semantic unit")
            x, y, w, h = image_box(item, asset_dir)
            slide.shapes.add_picture(str(asset_dir / item["path"]), emu(x, scale), emu(y, scale), emu(w, scale), emu(h, scale))
        elif typ == "line":
            add_line(slide, item, scale)
        elif typ in {"rect", "round_rect", "right_arrow", "down_arrow"}:
            add_shape(slide, item, scale)
        else:
            raise ValueError(f"Unsupported element type: {typ}")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build strict image-first PPTX replica.")
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--asset-dir", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    args = parser.parse_args()
    build(args.manifest, args.out, args.asset_dir)
    print(args.out)


if __name__ == "__main__":
    main()
